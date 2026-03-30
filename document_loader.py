"""
文档加载器 - 支持多种格式的文档加载
"""

import os
import json
import hashlib
import zipfile
from pathlib import Path
from typing import List
import pandas as pd
from datetime import datetime
from langchain_community.document_loaders import (
    TextLoader,
    UnstructuredMarkdownLoader,
    PyPDFLoader,
    CSVLoader,
    JSONLoader,
    UnstructuredWordDocumentLoader,
    UnstructuredExcelLoader,
    UnstructuredPowerPointLoader,
)
from langchain_core.documents import Document


def compute_file_hash(file_path: Path) -> str:
    """计算文件的MD5哈希值（基于完整内容）"""
    hasher = hashlib.md5()
    try:
        with open(file_path, 'rb') as f:
            for chunk in iter(lambda: f.read(8192), b''):
                hasher.update(chunk)
        return hasher.hexdigest()
    except Exception as e:
        print(f"计算文件哈希失败 {file_path.name}: {e}")
        return ""


class DocumentLoader:
    """文档加载器"""

    def __init__(self, docs_dir: Path):
        self.docs_dir = Path(docs_dir)
        if not self.docs_dir.exists():
            self.docs_dir.mkdir(parents=True, exist_ok=True)

    def _add_version_metadata(self, documents: List[Document], file_path: Path) -> List[Document]:
        """为文档添加版本信息到metadata"""
        file_hash = compute_file_hash(file_path)
        version = datetime.now().strftime("%Y%m%d_%H%M%S")

        for doc in documents:
            # 更新metadata，添加版本信息
            doc.metadata['file_hash'] = file_hash
            doc.metadata['version'] = version
            # 确保source字段存在
            if 'source' not in doc.metadata:
                doc.metadata['source'] = str(file_path)

        return documents

    def load_file(self, file_path: Path):
        """加载单个文件"""
        suffix = file_path.suffix.lower()

        # 特殊处理：跳过 .ppt (稳定性问题)
        if suffix == '.ppt':
            print(f"跳过(需转换格式): {file_path.name}")
            return None

        # 自定义加载器处理
        if suffix == '.xmind':
            return self._load_xmind(file_path)
        elif suffix == '.vsdx':
            return self._load_vsdx(file_path)
        elif suffix == '.docx':
            return self._load_docx(file_path)
        elif suffix == '.doc':
            return self._load_doc(file_path)
        elif suffix in ['.xlsx', '.xls']:
            return self._load_excel(file_path)
        elif suffix in ['.pptx', '.ppt']:
            return self._load_pptx(file_path)

        # 基础加载器
        loaders = {
            '.txt': TextLoader,
            '.md': UnstructuredMarkdownLoader,
            '.markdown': UnstructuredMarkdownLoader,
            '.pdf': PyPDFLoader,
            '.csv': CSVLoader,
            '.json': JSONLoader,
        }

        loader_class = loaders.get(suffix)
        if not loader_class:
            print(f"不支持的文件类型: {suffix}")
            return None

        try:
            if suffix == '.json':
                loader = loader_class(str(file_path), jq_schema='.')
            elif suffix == '.pdf':
                loader = loader_class(str(file_path))
            elif suffix == '.xlsx' or suffix == '.xls':
                loader = loader_class(str(file_path))
            elif suffix == '.csv':
                # CSV 文件尝试多种编码
                documents = self._load_csv_with_encoding(file_path)
                if documents:
                    documents = self._add_version_metadata(documents, file_path)
                    return documents
                return None
            else:
                loader = loader_class(str(file_path), encoding='utf-8')

            documents = loader.load()
            print(f"成功加载: {file_path.name}, 共 {len(documents)} 个文档")
            # 添加版本信息到metadata
            documents = self._add_version_metadata(documents, file_path)
            return documents
        except Exception as e:
            print(f"加载文件失败 {file_path.name}: {e}")
            return None

    def _load_csv_with_encoding(self, file_path: Path) -> List[Document]:
        """尝试多种编码加载 CSV 文件"""
        from langchain_community.document_loaders import CSVLoader

        # 尝试的编码顺序
        encodings = ['utf-8', 'gbk', 'gb2312', 'gb18030', 'utf-8-sig', 'latin1']

        for encoding in encodings:
            try:
                loader = CSVLoader(str(file_path), encoding=encoding)
                documents = loader.load()
                print(f"成功加载: {file_path.name}, 共 {len(documents)} 个文档 (编码: {encoding})")
                return documents
            except Exception as e:
                continue

        print(f"CSV文件所有编码尝试失败: {file_path.name}")
        return None

    def _load_docx(self, file_path: Path) -> List[Document]:
        """加载 Word .docx 文件"""
        try:
            from docx import Document as DocxDocument

            docx = DocxDocument(str(file_path))
            content_parts = []
            table_count = 0

            # 遍历文档元素，按顺序处理段落和表格
            from docx.oxml.table import CT_Tbl
            from docx.oxml.text.paragraph import CT_P
            from docx.table import Table
            from docx.text.paragraph import Paragraph

            body = docx._body._body
            for child in body:
                if isinstance(child, CT_P):
                    # 段落
                    para = Paragraph(child, docx)
                    if para.text.strip():
                        content_parts.append({
                            'type': 'paragraph',
                            'content': para.text.strip()
                        })
                elif isinstance(child, CT_Tbl):
                    # 表格
                    table = Table(child, docx)
                    md_table = self._convert_table_to_markdown(table)
                    if md_table:
                        table_count += 1
                        content_parts.append({
                            'type': 'table',
                            'content': md_table,
                            'table_header': self._extract_table_header(table)
                        })

            if content_parts:
                # 构建文档内容
                doc_contents = []
                for part in content_parts:
                    if part['type'] == 'paragraph':
                        doc_contents.append(part['content'])
                    elif part['type'] == 'table':
                        doc_contents.append(part['content'])

                content = '\n\n'.join(doc_contents)

                # 收集所有表格的表头元数据
                tables_metadata = [
                    p['table_header']
                    for p in content_parts if p['type'] == 'table'
                ]

                doc_metadata = {
                    'source': str(file_path),
                    'type': 'docx',
                    'table_count': table_count,
                    'has_tables': table_count > 0,
                }
                if tables_metadata:
                    doc_metadata['tables_metadata'] = tables_metadata
                print(f"成功加载: {file_path.name}, 包含 {table_count} 个表格")
                doc = Document(page_content=content, metadata=doc_metadata)
                return self._add_version_metadata([doc], file_path)
            else:
                print(f"Word文件内容为空: {file_path.name}")
                return None
        except Exception as e:
            print(f"加载Word文件失败 {file_path.name}: {e}")
            return None

    def _convert_table_to_markdown(self, table) -> str:
        """将Word表格转换为Markdown格式，正确处理合并单元格"""
        if not table.rows:
            return ""

        # 获取表格的实际列数（考虑合并单元格）
        num_cols = self._get_table_num_cols(table)
        if num_cols == 0:
            return ""

        # 构建单元格网格
        # grid[row][col] = text (None表示该位置被合并)
        grid = [[None] * num_cols for _ in range(len(table.rows))]

        # 处理每一行
        for row_idx, row in enumerate(table.rows):
            col_idx = 0
            for cell in row.cells:
                # 跳过已被水平或垂直合并占据的位置
                while col_idx < num_cols and grid[row_idx][col_idx] is not None:
                    col_idx += 1

                if col_idx >= num_cols:
                    break

                # 获取水平合并
                colspan = self._get_cell_colspan(cell)
                text = cell.text.strip().replace('\n', ' ')

                # 填充水平合并的单元格
                for c in range(colspan):
                    if col_idx + c < num_cols:
                        if c == 0:
                            grid[row_idx][col_idx + c] = text
                        else:
                            # 水平合并的占位单元格标记为''
                            grid[row_idx][col_idx + c] = ''

                col_idx += colspan

        # 检测并跳过标题行
        # 情况1：第一行跨列合并，colspan>1
        # 情况2：第一行所有非空单元格内容相同（Word用重复文本代替合并）
        data_start_row = 0
        if len(table.rows) > 1:
            first_row_texts = [grid[0][c] for c in range(num_cols) if grid[0][c]]
            if first_row_texts:
                unique_texts = set(first_row_texts)
                # 检查第一行是否有水平合并
                has_colspan = any(self._get_cell_colspan(cell) > 1 for cell in table.rows[0].cells)
                # 情况1：colspan>1
                if has_colspan:
                    data_start_row = 1
                # 情况2：所有非空内容相同，且数量少于列数
                elif len(unique_texts) == 1 and len(first_row_texts) < num_cols:
                    data_start_row = 1
                # 情况3：只有一个单元格有内容，其他都是空
                elif len(unique_texts) == 1 and len(first_row_texts) == 1 and num_cols > 1:
                    data_start_row = 1
                # 情况4：所有非空内容完全相同，且列数>1
                elif len(unique_texts) == 1 and num_cols > 1:
                    data_start_row = 1

        # 生成Markdown表格
        lines = []

        # 表头行
        header_row = [str(grid[data_start_row][c]) if grid[data_start_row][c] is not None else '' for c in range(num_cols)]
        lines.append('| ' + ' | '.join(header_row) + ' |')

        # 分隔行
        lines.append('|' + '---|' * num_cols)

        # 数据行
        for row_idx in range(data_start_row + 1, len(table.rows)):
            row_data = [str(grid[row_idx][c]) if grid[row_idx][c] is not None else '' for c in range(num_cols)]
            lines.append('| ' + ' | '.join(row_data) + ' |')

        return '\n'.join(lines)

    def _get_table_num_cols(self, table) -> int:
        """计算表格的实际列数（考虑合并单元格）"""
        max_cols = 0
        for row in table.rows:
            col_count = 0
            for cell in row.cells:
                col_count += self._get_cell_colspan(cell)
            max_cols = max(max_cols, col_count)
        return max_cols

    def _get_cell_colspan(self, cell) -> int:
        """获取单元格的水平合并数（colspan）"""
        try:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            gridSpan = tcPr.find('w:gridSpan')
            if gridSpan is not None:
                val = gridSpan.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val')
                return int(val) if val else 1
        except Exception:
            pass
        return 1

    def _extract_table_header(self, table) -> str:
        """提取表格列头（跳过后检测后的实际数据列头）"""
        if not table.rows:
            return ""

        # 计算实际列数
        num_cols = self._get_table_num_cols(table)
        if num_cols == 0:
            return ""

        # 构建grid
        grid = [[None] * num_cols for _ in range(len(table.rows))]
        for row_idx, row in enumerate(table.rows):
            col_idx = 0
            for cell in row.cells:
                while col_idx < num_cols and grid[row_idx][col_idx] is not None:
                    col_idx += 1
                if col_idx >= num_cols:
                    break
                colspan = self._get_cell_colspan(cell)
                text = cell.text.strip().replace('\n', ' ')[:20]
                for c in range(colspan):
                    if col_idx + c < num_cols:
                        grid[row_idx][col_idx + c] = text if c == 0 else ''
                col_idx += colspan

        # 检测并跳过标题行
        data_start_row = 0
        if len(table.rows) > 1:
            first_row_texts = [grid[0][c] for c in range(num_cols) if grid[0][c]]
            if first_row_texts:
                unique_texts = set(first_row_texts)
                has_colspan = any(self._get_cell_colspan(cell) > 1 for cell in table.rows[0].cells)
                if has_colspan or (len(unique_texts) == 1 and num_cols > 1):
                    data_start_row = 1

        # 提取实际的列头
        if data_start_row < len(table.rows):
            header_cells = [str(grid[data_start_row][c]) if grid[data_start_row][c] else '' for c in range(num_cols)]
            return ' | '.join([h for h in header_cells if h][:5])
        return ''

    def _load_doc(self, file_path: Path) -> List[Document]:
        """加载 Word .doc 文件（旧格式）"""
        try:
            # 尝试使用 pywin32 读取 .doc 文件
            import win32com.client
            import os

            # 创建 Word 应用实例
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False

            try:
                # 打开文档
                doc = word.Documents.Open(str(file_path.absolute()))
                paragraphs = []

                # 提取所有段落
                for para in doc.Paragraphs:
                    text = para.Range.Text.strip()
                    if text:
                        paragraphs.append(text)

                # 释放文档
                doc.Close(False)

                if paragraphs:
                    content = '\n\n'.join(paragraphs)
                    doc = Document(
                        page_content=content,
                        metadata={'source': str(file_path), 'type': 'doc'}
                    )
                    print(f"成功加载: {file_path.name}")
                    return self._add_version_metadata([doc], file_path)
                else:
                    print(f"Word文件内容为空: {file_path.name}")
                    return None
            finally:
                word.Quit()
        except ImportError:
            print(f"读取 .doc 文件需要 pywin32 库，请运行: pip install pywin32")
            print(f"或者将 {file_path.name} 转换为 .docx 格式")
            return None
        except Exception as e:
            print(f"加载Word .doc文件失败 {file_path.name}: {e}")
            return None

    def _load_excel(self, file_path: Path) -> List[Document]:
        """加载 Excel 文件，转换为Markdown表格"""
        try:
            import pandas as pd

            content_parts = []
            table_count = 0

            # 读取所有sheet
            excel_file = pd.ExcelFile(str(file_path))
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name, nrows=500)

                if df.empty:
                    continue

                # 转换为Markdown表格
                md_table = self._convert_dataframe_to_markdown(df, sheet_name)
                if md_table:
                    table_count += 1
                    content_parts.append({
                        'type': 'table',
                        'content': md_table,
                        'table_header': sheet_name,
                        'sheet_name': sheet_name
                    })

            if content_parts:
                # 构建文档内容
                doc_contents = []
                tables_metadata = []
                for part in content_parts:
                    doc_contents.append(part['content'])
                    tables_metadata.append(f"{part.get('sheet_name', '')}: {part['table_header']}")

                content = '\n\n'.join(doc_contents)

                doc_metadata = {
                    'source': str(file_path),
                    'type': 'excel',
                    'table_count': table_count,
                    'has_tables': table_count > 0,
                }
                if tables_metadata:
                    doc_metadata['tables_metadata'] = tables_metadata

                doc = Document(page_content=content, metadata=doc_metadata)
                print(f"成功加载: {file_path.name}, 包含 {table_count} 个表格")
                # 添加版本信息
                return self._add_version_metadata([doc], file_path)
            else:
                print(f"Excel文件内容为空: {file_path.name}")
                return None
        except Exception as e:
            print(f"加载Excel文件失败 {file_path.name}: {e}")
            return None

    def _convert_dataframe_to_markdown(self, df: pd.DataFrame, sheet_name: str = "") -> str:
        """将DataFrame转换为Markdown表格"""
        if df.empty:
            return ""

        lines = []

        # 清理列名
        headers = [str(col)[:50] for col in df.columns]
        lines.append('| ' + ' | '.join(headers) + ' |')
        lines.append('|' + '---|' * len(headers))

        # 添加数据行
        for _, row in df.iterrows():
            cells = [str(val)[:100] if pd.notna(val) else ' ' for val in row]
            lines.append('| ' + ' | '.join(cells) + ' |')

        return '\n'.join(lines)

    def _load_pptx(self, file_path: Path) -> List[Document]:
        """加载 PowerPoint .pptx 文件"""
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            paragraphs = []

            # 提取所有幻灯片的文本
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"\n--- Slide {slide_num} ---\n"]

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                slide_text.append(para.text)

                    # 处理表格
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text)
                            if row_text:
                                slide_text.append(' | '.join(row_text))

                if len(slide_text) > 1:  # 除了标题还有内容
                    paragraphs.append(''.join(slide_text))

            if paragraphs:
                content = '\n\n'.join(paragraphs)
                doc = Document(
                    page_content=content,
                    metadata={'source': str(file_path), 'type': 'pptx'}
                )
                print(f"成功加载: {file_path.name}")
                # 添加版本信息
                return self._add_version_metadata([doc], file_path)
            else:
                print(f"PPT文件内容为空: {file_path.name}")
                return None
        except Exception as e:
            print(f"加载PPT文件失败 {file_path.name}: {e}")
            return None

    def _load_xmind(self, file_path: Path) -> List[Document]:
        """加载 XMind 思维导图文件"""
        try:
            from xmindparser import xmind_to_dict

            # xmind_to_dict 返回结构: {'topics': [...], 'title': '...'}
            data = xmind_to_dict(str(file_path))

            # 转换为文本
            text = self._parse_xmind_dict(data)

            if text:
                doc = Document(
                    page_content=text,
                    metadata={'source': str(file_path), 'type': 'xmind'}
                )
                print(f"成功加载: {file_path.name}")
                # 添加版本信息
                return self._add_version_metadata([doc], file_path)
            else:
                print(f"XMind文件内容为空: {file_path.name}")
                return None
        except Exception as e:
            print(f"加载XMind文件失败 {file_path.name}: {e}")
            return None

    def _parse_xmind_dict(self, data, level=0) -> str:
        """递归解析 XMind 数据为文本"""
        lines = []

        if isinstance(data, dict):
            # 获取标题
            title = data.get('title', '')
            if title:
                prefix = '#' * min(level + 1, 6)
                lines.append(f"{prefix} {title}")

            # 处理子主题
            topics = data.get('topics', [])
            if isinstance(topics, dict):
                topics = [topics]
            for topic in topics:
                if isinstance(topic, dict):
                    lines.append(self._parse_xmind_dict(topic, level + 1))

        elif isinstance(data, list):
            for item in data:
                lines.append(self._parse_xmind_dict(item, level))

        return '\n'.join([l for l in lines if l])

    def _load_vsdx(self, file_path: Path) -> List[Document]:
        """加载 Visio vsdx 文件"""
        try:
            content = []
            with zipfile.ZipFile(file_path, 'r') as zf:
                # vsdx 文件中的内容在 document.xml 中
                if 'document.xml' in zf.namelist():
                    with zf.open('document.xml') as f:
                        xml_content = f.read().decode('utf-8')
                        text = self._extract_text_from_xml(xml_content)
                        if text:
                            content.append(text)

                # 页面信息
                if 'documentRelationship.xml' in zf.namelist():
                    pass  # 可以扩展处理页面信息

            if content:
                text = '\n\n'.join(content)
                doc = Document(
                    page_content=text,
                    metadata={'source': str(file_path), 'type': 'vsdx'}
                )
                print(f"成功加载: {file_path.name}")
                # 添加版本信息
                return self._add_version_metadata([doc], file_path)
            else:
                print(f"Visio文件内容为空: {file_path.name}")
                return None
        except Exception as e:
            print(f"加载Visio文件失败 {file_path.name}: {e}")
            return None

    def _extract_text_from_xml(self, xml_content: str) -> str:
        """从 XML 中提取文本"""
        import re
        # 提取 <a:t> 标签中的文本（Office Open XML 格式）
        texts = re.findall(r'<a:t[^>]*>([^<]*)</a:t>', xml_content, re.IGNORECASE)
        return '\n'.join(texts)

    def load_directory(self) -> List:
        """加载目录下的所有文档"""
        all_docs = []
        supported = set(get_supported_extensions())

        # 获取所有文件
        all_files = list(self.docs_dir.rglob('*'))
        total = len([f for f in all_files if f.is_file()])
        processed = 0

        print(f"开始加载文档，共 {total} 个文件...")

        for file_path in all_files:
            if not file_path.is_file():
                continue

            # 跳过不支持的文件类型
            if file_path.suffix.lower() not in supported:
                continue

            processed += 1
            print(f"[{processed}/{total}] 加载: {file_path.name}")

            try:
                docs = self.load_file(file_path)
                if docs:
                    all_docs.extend(docs)
            except Exception as e:
                print(f"加载失败 {file_path.name}: {e}")

        print(f"文档加载完成，共 {len(all_docs)} 个文档对象")
        return all_docs


def get_supported_extensions():
    """获取支持的文件扩展名"""
    return ['.txt', '.md', '.markdown', '.pdf', '.csv', '.json', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.xmind', '.vsdx']
